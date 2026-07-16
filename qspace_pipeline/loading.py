"""Component 2: robust multi-format document loading."""
from __future__ import annotations

import csv as csvmod
import io
import json
import re
from html.parser import HTMLParser
from pathlib import Path
from typing import Dict, List, Tuple

from .config import Paths
from .utils import banner, estimate_tokens, log, sub, truncate, write_jsonl


class _HTMLText(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._parts: List[str] = []
        self._skip = False

    def handle_starttag(self, tag, attrs):
        if tag in ("script", "style"):
            self._skip = True

    def handle_endtag(self, tag):
        if tag in ("script", "style"):
            self._skip = False

    def handle_data(self, data):
        if not self._skip and data.strip():
            self._parts.append(data.strip())

    def text(self) -> str:
        return "\n".join(self._parts)


def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _load_markdown(path: Path) -> Tuple[str, Dict, List[str]]:
    text = _read_text(path)
    headings = [ln.strip("# ").strip() for ln in text.splitlines()
                if ln.lstrip().startswith("#")]
    return text, {"headings": headings}, []


def _load_text(path: Path) -> Tuple[str, Dict, List[str]]:
    return _read_text(path), {}, []


def _load_html(path: Path) -> Tuple[str, Dict, List[str]]:
    parser = _HTMLText()
    try:
        parser.feed(_read_text(path))
        return parser.text(), {}, []
    except Exception as exc:
        return _read_text(path), {}, [f"html-parse-fallback: {exc}"]


def _load_csv(path: Path) -> Tuple[str, Dict, List[str]]:
    warnings: List[str] = []
    rows: List[List[str]] = []
    try:
        with path.open("r", encoding="utf-8", errors="ignore", newline="") as fh:
            reader = csvmod.reader(fh)
            rows = [r for r in reader]
    except Exception as exc:
        return _read_text(path), {}, [f"csv-parse-fallback: {exc}"]
    if not rows:
        return "", {"tables": []}, ["empty-csv"]
    header = rows[0]
    # Render as readable sentences so answer detection can work on structured data.
    lines: List[str] = [f"Table columns: {', '.join(header)}."]
    table_records: List[Dict[str, str]] = []
    for r in rows[1:]:
        record = {header[i] if i < len(header) else f"col{i}": v
                  for i, v in enumerate(r)}
        table_records.append(record)
        parts = [f"{k} = {v}" for k, v in record.items() if v != ""]
        if parts:
            lines.append("Row: " + "; ".join(parts) + ".")
    meta = {"tables": [{"columns": header, "records": table_records}]}
    return "\n".join(lines), meta, warnings


def _flatten_json(obj, prefix: str = "") -> List[str]:
    lines: List[str] = []
    if isinstance(obj, dict):
        for k, v in obj.items():
            if isinstance(v, (dict, list)):
                lines += _flatten_json(v, f"{prefix}{k}.")
            else:
                lines.append(f"{prefix}{k} = {v}")
    elif isinstance(obj, list):
        for i, v in enumerate(obj):
            if isinstance(v, (dict, list)):
                lines += _flatten_json(v, f"{prefix}[{i}].")
            else:
                lines.append(f"{prefix}[{i}] = {v}")
    else:
        lines.append(f"{prefix} = {obj}")
    return lines


def _load_json(path: Path) -> Tuple[str, Dict, List[str]]:
    try:
        obj = json.loads(_read_text(path))
    except Exception as exc:
        return _read_text(path), {}, [f"json-parse-fallback: {exc}"]
    lines = _flatten_json(obj)
    text = "\n".join(lines)
    return text, {"json_keys": list(obj.keys()) if isinstance(obj, dict) else []}, []


def _load_jsonl(path: Path) -> Tuple[str, Dict, List[str]]:
    lines: List[str] = []
    warnings: List[str] = []
    for i, raw in enumerate(_read_text(path).splitlines()):
        raw = raw.strip()
        if not raw:
            continue
        try:
            obj = json.loads(raw)
            lines += _flatten_json(obj, f"record[{i}].")
        except Exception as exc:
            warnings.append(f"line {i} parse error: {exc}")
    return "\n".join(lines), {}, warnings


def _load_pdf(path: Path) -> Tuple[str, Dict, List[str]]:
    for mod_name, extractor in (("pypdf", "pypdf"), ("PyPDF2", "pypdf2"),
                                ("pdfplumber", "pdfplumber")):
        try:
            if extractor in ("pypdf", "pypdf2"):
                mod = __import__(mod_name)
                reader = mod.PdfReader(str(path))
                text = "\n".join((pg.extract_text() or "") for pg in reader.pages)
                return text, {"pages": len(reader.pages)}, []
            else:
                import pdfplumber  # type: ignore
                with pdfplumber.open(str(path)) as pdf:
                    text = "\n".join((pg.extract_text() or "") for pg in pdf.pages)
                    return text, {"pages": len(pdf.pages)}, []
        except Exception:
            continue
    return "", {}, ["no pdf loader available (install pypdf); skipped extraction"]


def _load_docx(path: Path) -> Tuple[str, Dict, List[str]]:
    try:
        import docx  # type: ignore
        d = docx.Document(str(path))
        text = "\n".join(p.text for p in d.paragraphs)
        return text, {}, []
    except Exception:
        return "", {}, ["no docx loader available (install python-docx); skipped"]


def _load_xlsx(path: Path) -> Tuple[str, Dict, List[str]]:
    try:
        import openpyxl  # type: ignore
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        lines: List[str] = []
        for ws in wb.worksheets:
            lines.append(f"Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    lines.append("Row: " + "; ".join(cells))
        return "\n".join(lines), {}, []
    except Exception:
        return "", {}, ["no xlsx loader available (install openpyxl); skipped"]


def _load_pptx(path: Path) -> Tuple[str, Dict, List[str]]:
    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation(str(path))
        lines: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    lines.append(shape.text_frame.text)
        return "\n".join(lines), {}, []
    except Exception:
        return "", {}, ["no pptx loader available (install python-pptx); skipped"]


_DISPATCH = {
    "markdown": _load_markdown, "text": _load_text, "html": _load_html,
    "csv": _load_csv, "json": _load_json, "jsonl": _load_jsonl,
    "pdf": _load_pdf, "docx": _load_docx, "xlsx": _load_xlsx, "pptx": _load_pptx,
}


def load_documents(paths: Paths, inventory: List[Dict]) -> List[Dict]:
    banner("LOAD", "Document loading")
    loaded: List[Dict] = []
    ok = 0
    failed: List[str] = []
    for row in inventory:
        loader = row["detected_loader"]
        path = Path(row["file_path"])
        fn = _DISPATCH.get(loader)
        warnings: List[str] = []
        text, meta = "", {}
        if fn is None:
            warnings = [f"unsupported loader for {row['extension']}"]
            status = "unsupported"
        else:
            try:
                text, meta, warnings = fn(path)
                status = "loaded" if text.strip() else "empty"
            except Exception as exc:  # never crash the pipeline
                status = "error"
                warnings = [f"loader raised: {exc}"]
        row["load_status"] = status
        row["character_count"] = len(text)
        row["token_estimate"] = estimate_tokens(text)
        if status == "loaded":
            ok += 1
        else:
            failed.append(f"{row['file_name']} ({status})")
        loaded.append({
            "document_id": row["document_id"],
            "file_path": row["file_path"],
            "title": row["document_title_guess"],
            "extension": row["extension"],
            "domain_hints": row["possible_domain_hints"],
            "text": text,
            "metadata": meta,
            "load_status": status,
            "load_warnings": warnings,
        })

    write_jsonl(paths.intermediate / "loaded_documents.jsonl", loaded)

    log("LOAD", f"Loaded {ok} / {len(inventory)} documents with usable text.")
    if failed:
        log("LOAD", f"Non-loaded / empty files: {len(failed)}")
        for f in failed[:10]:
            sub(f"- {f}")
    log("LOAD", "Sample loaded snippets:")
    shown = [d for d in loaded if d["load_status"] == "loaded"][:3]
    for d in shown:
        sub(f"{d['document_id']} [{d['title']}]: "
            f"\"{truncate(d['text'], 110)}\"")
    log("LOAD", f"Saved: {paths.intermediate / 'loaded_documents.jsonl'}")
    return loaded
